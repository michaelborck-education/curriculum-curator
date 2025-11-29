"""
File import service for processing uploaded course materials
Supports PDF, DOCX, PPTX, and other document formats
"""

import io
import mimetypes
import re
from collections import Counter
from pathlib import Path
from typing import Any, ClassVar

# Document processing libraries - optional dependencies
PyPDF2: Any = None
Document: Any = None
Presentation: Any = None
markdown: Any = None

try:
    import PyPDF2

    has_pdf = True
except ImportError:
    has_pdf = False

try:
    from docx import Document

    has_docx = True
except ImportError:
    has_docx = False

try:
    from pptx import Presentation

    has_pptx = True
except ImportError:
    has_pptx = False

try:
    import markdown

    has_markdown = True
except ImportError:
    has_markdown = False


class FileImportService:
    """Service for importing and processing various file formats"""

    SUPPORTED_FORMATS: ClassVar[dict[str, list[str]]] = {
        "application/pdf": [".pdf"],
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": [
            ".docx"
        ],
        "application/msword": [".doc"],
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": [
            ".pptx"
        ],
        "application/vnd.ms-powerpoint": [".ppt"],
        "text/markdown": [".md"],
        "text/plain": [".txt"],
        "text/html": [".html", ".htm"],
    }

    def __init__(self):
        # Extended content types with more patterns
        self.content_patterns = {
            "lecture": [
                r"lecture\s*\d+",
                r"chapter\s*\d+",
                r"module\s*\d+",
                r"lesson\s*\d+",
                r"topic:",
                r"learning objectives?:",
                r"introduction:",
                r"overview:",
                r"today we will",
                r"in this lecture",
            ],
            "quiz": [
                r"quiz\s*\d*",
                r"test\s*\d*",
                r"exam\s*\d*",
                r"assessment",
                r"question\s*\d+",
                r"q\d+\.",
                r"multiple choice",
                r"true/false",
                r"select the correct",
                r"fill in the blank",
            ],
            "worksheet": [
                r"worksheet",
                r"exercise\s*\d*",
                r"practice problems?",
                r"homework",
                r"assignment\s*\d*",
                r"problem set",
                r"complete the following",
                r"work through",
            ],
            "lab": [
                r"lab\s*\d*",
                r"laboratory",
                r"experiment\s*\d*",
                r"practical\s*\d*",
                r"hands-on",
                r"procedure:",
                r"materials:",
                r"equipment:",
                r"apparatus",
            ],
            "case_study": [
                r"case study",
                r"case\s*\d+",
                r"scenario:",
                r"situation:",
                r"background:",
                r"analysis:",
                r"discussion questions?",
                r"real-world example",
            ],
            "interactive": [
                r"interactive",
                r"simulation",
                r"game",
                r"click",
                r"drag",
                r"select",
                r"<button",
                r"<input",
                r"javascript:",
                r"onclick",
            ],
            "presentation": [
                r"presentation",
                r"slides?",
                r"slide\s*\d+",
                r"agenda:",
                r"outline:",
                r"summary slide",
            ],
            "reading": [
                r"reading",
                r"article",
                r"paper",
                r"journal",
                r"abstract:",
                r"references:",
                r"bibliography",
            ],
            "video_script": [
                r"video",
                r"transcript",
                r"narration",
                r"voiceover",
                r"[music]",
                r"[scene",
                r"fade in",
                r"cut to",
            ],
        }

        # Content type metadata
        self.content_type_info = {
            "lecture": {"name": "Lecture", "icon": "📚", "default_duration": 45},
            "quiz": {"name": "Quiz/Assessment", "icon": "📝", "default_duration": 30},
            "worksheet": {
                "name": "Worksheet/Exercise",
                "icon": "✏️",
                "default_duration": 30,
            },
            "lab": {"name": "Lab/Practical", "icon": "🔬", "default_duration": 90},
            "case_study": {"name": "Case Study", "icon": "💼", "default_duration": 45},
            "interactive": {
                "name": "Interactive Content",
                "icon": "🎮",
                "default_duration": 30,
            },
            "presentation": {
                "name": "Presentation",
                "icon": "🎯",
                "default_duration": 30,
            },
            "reading": {
                "name": "Reading Material",
                "icon": "📖",
                "default_duration": 20,
            },
            "video_script": {
                "name": "Video/Media",
                "icon": "🎥",
                "default_duration": 15,
            },
            "general": {
                "name": "General Content",
                "icon": "📄",
                "default_duration": 30,
            },
        }

    async def process_file(
        self, file_content: bytes, filename: str, content_type: str | None = None
    ) -> dict[str, Any]:
        """
        Process an uploaded file and extract its content

        Returns:
            Dict containing:
            - content: Extracted text content
            - content_type: Detected type (lecture, quiz, etc.)
            - metadata: File metadata
            - sections: Parsed sections/structure
            - suggestions: Enhancement suggestions
        """
        # Determine file type
        if not content_type:
            content_type, _ = mimetypes.guess_type(filename)

        file_extension = Path(filename).suffix.lower()

        # Extract content based on file type
        extracted_content = ""
        metadata = {
            "filename": filename,
            "size": len(file_content),
            "type": content_type,
            "extension": file_extension,
        }

        try:
            if file_extension == ".pdf" and has_pdf:
                extracted_content = await self._extract_pdf(file_content)
            elif file_extension == ".docx" and has_docx:
                extracted_content = await self._extract_docx(file_content)
            elif file_extension in [".ppt", ".pptx"] and has_pptx:
                extracted_content = await self._extract_pptx(file_content)
            elif file_extension == ".md" and has_markdown:
                extracted_content = await self._extract_markdown(file_content)
            elif file_extension in [".txt", ".html", ".htm"]:
                extracted_content = file_content.decode("utf-8", errors="ignore")
            else:
                raise ValueError(f"Unsupported file format: {file_extension}")

            # Detect content type with confidence
            detected_type, confidence, type_scores = self._detect_content_type(
                extracted_content
            )

            # Parse sections
            sections = self._parse_sections(extracted_content)

            # Generate enhancement suggestions
            suggestions = self._generate_suggestions(extracted_content, detected_type)

            # Perform gap analysis
            gaps = self._analyze_gaps(extracted_content, detected_type)

            # Get all available content types for user selection
            available_types = list(self.content_type_info.keys())

            return {
                "success": True,
                "content": extracted_content,
                "content_type": detected_type,
                "content_type_confidence": confidence,
                "content_type_scores": type_scores,
                "available_content_types": available_types,
                "content_type_info": self.content_type_info,
                "metadata": metadata,
                "sections": sections,
                "suggestions": suggestions,
                "gaps": gaps,
                "word_count": len(extracted_content.split()),
                "estimated_reading_time": len(extracted_content.split())
                // 200,  # minutes
            }

        except Exception as e:
            return {
                "success": False,
                "error": str(e),
                "metadata": metadata,
                "content": "",
                "content_type": "unknown",
                "sections": [],
                "suggestions": [],
                "gaps": [],
            }

    async def _extract_pdf(self, file_content: bytes) -> str:
        """Extract text from PDF file"""
        if not has_pdf:
            raise ValueError("PDF processing not available. Install PyPDF2.")

        text_content = []
        pdf_file = io.BytesIO(file_content)
        pdf_reader = PyPDF2.PdfReader(pdf_file)

        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            text = page.extract_text()
            if text:
                text_content.append(f"--- Page {page_num + 1} ---\n{text}")

        return "\n\n".join(text_content)

    async def _extract_docx(self, file_content: bytes) -> str:
        """Extract text from DOCX file"""
        if not has_docx:
            raise ValueError("DOCX processing not available. Install python-docx.")

        text_content = []
        docx_file = io.BytesIO(file_content)
        doc = Document(docx_file)

        text_content.extend(
            paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()
        )

        # Also extract from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = [
                    cell.text.strip() for cell in row.cells if cell.text.strip()
                ]
                if row_text:
                    text_content.append(" | ".join(row_text))

        return "\n\n".join(text_content)

    async def _extract_pptx(self, file_content: bytes) -> str:
        """Extract text from PowerPoint file"""
        if not has_pptx:
            raise ValueError("PPTX processing not available. Install python-pptx.")

        text_content = []
        pptx_file = io.BytesIO(file_content)
        presentation = Presentation(pptx_file)

        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_text = [f"--- Slide {slide_num} ---"]

            # Extract title
            if slide.shapes.title:
                slide_text.append(f"Title: {slide.shapes.title.text}")

            # Extract content from all shapes (avoiding duplicate title)
            slide_text.extend(
                shape.text
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text and shape != slide.shapes.title
            )

            # Extract notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text:
                slide_text.append(f"Notes: {slide.notes_slide.notes_text_frame.text}")

            if len(slide_text) > 1:  # Only add if there's content beyond the header
                text_content.append("\n".join(slide_text))

        return "\n\n".join(text_content)

    async def _extract_markdown(self, file_content: bytes) -> str:
        """Extract and process markdown content"""
        if not has_markdown:
            # Just return raw text if markdown not available
            return file_content.decode("utf-8", errors="ignore")

        text = file_content.decode("utf-8", errors="ignore")
        # Convert markdown to plain text while preserving structure
        html = markdown.markdown(text)
        # Simple HTML to text conversion
        return re.sub(r"<[^>]+>", "", html)

    def _detect_content_type(self, content: str) -> tuple[str, float, dict]:
        """
        Detect the type of content based on patterns

        Returns:
            tuple: (detected_type, confidence, scores_dict)
        """
        content_lower = content.lower()
        scores = {}

        for content_type, patterns in self.content_patterns.items():
            score = 0
            for pattern in patterns:
                matches = re.findall(pattern, content_lower)
                score += len(matches)
            scores[content_type] = score

        # Calculate confidence based on score distribution
        max_score = max(scores.values()) if scores else 0
        sum(scores.values())

        if max_score == 0:
            # No patterns matched at all
            return "general", 0.0, scores

        # Get the type with highest score
        detected_type = max(scores, key=lambda k: scores.get(k, 0))

        # Calculate confidence (0-1 scale)
        # High confidence if one type dominates, low if scores are close
        confidence = 0.0
        if max_score > 2:  # Minimum threshold
            # Confidence based on how much this type dominates
            second_highest = (
                sorted(scores.values(), reverse=True)[1] if len(scores) > 1 else 0
            )
            confidence = min(
                1.0, (max_score - second_highest) / max_score * (max_score / 10)
            )
            confidence = max(0.3, confidence)  # Minimum 30% if we detected something

        # If confidence is too low, return 'general'
        if confidence < 0.3:
            return "general", confidence, scores

        return detected_type, confidence, scores

    def _parse_sections(self, content: str) -> list[dict[str, str]]:
        """Parse content into logical sections"""
        sections = []

        # Common section headers
        section_patterns = [
            r"^#+\s+(.+)$",  # Markdown headers
            r"^([A-Z][^.!?]*):$",  # Title case with colon
            r"^(\d+\.?\s+[A-Z].+)$",  # Numbered sections
            r"^(Chapter|Section|Module|Unit|Lesson|Topic)\s+\d+:?\s*(.*)$",  # Common prefixes
            r"^(Introduction|Overview|Summary|Conclusion|References|Objectives|Materials|Procedure)$",  # Common sections
        ]

        lines = content.split("\n")
        current_section = {"title": "Introduction", "content": [], "start_line": 0}

        for i, line in enumerate(lines):
            stripped_line = line.strip()
            if not stripped_line:
                continue

            # Check if this line is a section header
            is_header = False
            for pattern in section_patterns:
                match = re.match(pattern, stripped_line, re.IGNORECASE)
                if match:
                    # Save current section if it has content
                    if current_section["content"]:
                        current_section["content"] = "\n".join(
                            current_section["content"]
                        )
                        sections.append(current_section)

                    # Start new section
                    title = (
                        match.group(1)
                        if match.lastindex and match.lastindex >= 1
                        else stripped_line
                    )
                    if (
                        match.lastindex and match.lastindex == 2
                    ):  # For patterns with two groups
                        title = f"{match.group(1)} {match.group(2)}".strip()

                    current_section = {"title": title, "content": [], "start_line": i}
                    is_header = True
                    break

            if not is_header:
                current_section["content"].append(stripped_line)

        # Add final section
        if current_section["content"]:
            current_section["content"] = "\n".join(current_section["content"])
            sections.append(current_section)

        return sections

    def _generate_suggestions(self, content: str, content_type: str) -> list[str]:
        """Generate enhancement suggestions based on content analysis"""
        suggestions = []
        content_lower = content.lower()

        # Add general suggestions
        suggestions.extend(self._get_general_suggestions(content, content_lower))

        # Add type-specific suggestions
        type_suggestions = self._get_type_specific_suggestions(
            content, content_lower, content_type
        )
        suggestions.extend(type_suggestions)

        # Add accessibility suggestions
        suggestions.extend(self._get_accessibility_suggestions(content_lower))

        return suggestions

    def _get_general_suggestions(self, content: str, content_lower: str) -> list[str]:
        """Get general content suggestions."""
        suggestions = []

        if len(content.split()) < 300:
            suggestions.append(
                "Content appears brief. Consider expanding with more details."
            )

        if "objective" not in content_lower and "goal" not in content_lower:
            suggestions.append("Add clear learning objectives at the beginning.")

        return suggestions

    def _get_type_specific_suggestions(
        self, content: str, content_lower: str, content_type: str
    ) -> list[str]:
        """Get content-type specific suggestions."""

        type_handlers = {
            "lecture": self._get_lecture_suggestions,
            "quiz": self._get_quiz_suggestions,
            "worksheet": self._get_worksheet_suggestions,
            "lab": self._get_lab_suggestions,
        }

        handler = type_handlers.get(content_type)
        if handler:
            return handler(content, content_lower)

        return []

    def _get_lecture_suggestions(self, content: str, content_lower: str) -> list[str]:
        """Get lecture-specific suggestions."""
        suggestions = []

        if "example" not in content_lower:
            suggestions.append("Include practical examples to illustrate concepts.")
        if "summary" not in content_lower and "conclusion" not in content_lower:
            suggestions.append("Add a summary or conclusion section.")
        if not re.search(r"\?", content):
            suggestions.append("Consider adding reflection questions for students.")

        return suggestions

    def _get_quiz_suggestions(self, content: str, content_lower: str) -> list[str]:
        """Get quiz-specific suggestions."""
        suggestions = []

        if "answer" not in content_lower and "solution" not in content_lower:
            suggestions.append("Include answer key or solutions.")
        if "point" not in content_lower and "score" not in content_lower:
            suggestions.append("Add point values for each question.")
        if "instruction" not in content_lower:
            suggestions.append("Add clear instructions for completing the assessment.")

        return suggestions

    def _get_worksheet_suggestions(self, content: str, content_lower: str) -> list[str]:
        """Get worksheet-specific suggestions."""
        suggestions = []

        if "instruction" not in content_lower:
            suggestions.append("Add clear instructions for each exercise.")
        if "example" not in content_lower:
            suggestions.append("Include worked examples to guide students.")

        return suggestions

    def _get_lab_suggestions(self, content: str, content_lower: str) -> list[str]:
        """Get lab-specific suggestions."""
        suggestions = []

        if "safety" not in content_lower:
            suggestions.append("Add safety guidelines and precautions.")
        if "material" not in content_lower and "equipment" not in content_lower:
            suggestions.append("Include complete list of required materials.")
        if "procedure" not in content_lower and "step" not in content_lower:
            suggestions.append("Add detailed step-by-step procedure.")

        return suggestions

    def _get_accessibility_suggestions(self, content_lower: str) -> list[str]:
        """Get accessibility-related suggestions."""
        suggestions = []

        if not re.search(r"(image|figure|diagram|chart|graph)", content_lower):
            suggestions.append("Consider adding visual aids to enhance understanding.")
        else:
            suggestions.append(
                "Ensure all images have descriptive alt text for accessibility."
            )

        return suggestions

    def _analyze_gaps(self, content: str, content_type: str) -> list[dict[str, str]]:
        """Analyze content for gaps and missing elements"""
        gaps = []
        content_lower = content.lower()

        # Required elements based on content type
        required_elements = {
            "lecture": {
                "learning_objectives": [
                    "objective",
                    "goal",
                    "outcome",
                    "will be able to",
                ],
                "introduction": ["introduction", "overview", "background"],
                "main_content": ["topic", "concept", "theory", "principle"],
                "examples": ["example", "for instance", "such as", "e.g."],
                "summary": ["summary", "conclusion", "recap", "key points"],
            },
            "quiz": {
                "instructions": ["instruction", "direction", "complete", "answer"],
                "questions": ["question", "q.", "problem"],
                "answer_format": [
                    "multiple choice",
                    "true/false",
                    "short answer",
                    "essay",
                ],
                "scoring": ["point", "score", "mark", "grade"],
            },
            "worksheet": {
                "instructions": ["instruction", "complete", "solve", "calculate"],
                "problems": ["problem", "exercise", "question", "task"],
                "workspace": ["space", "show your work", "solution area"],
            },
            "lab": {
                "objectives": ["objective", "purpose", "goal"],
                "materials": ["material", "equipment", "apparatus", "tool"],
                "procedure": ["procedure", "method", "step", "instruction"],
                "safety": ["safety", "caution", "warning", "hazard"],
                "data_collection": ["data", "observation", "record", "measure"],
                "analysis": ["analysis", "calculate", "interpret", "discuss"],
            },
        }

        elements = required_elements.get(content_type, required_elements["lecture"])

        for element_name, keywords in elements.items():
            found = any(keyword in content_lower for keyword in keywords)
            if not found:
                gaps.append(
                    {
                        "element": element_name.replace("_", " ").title(),
                        "severity": "high"
                        if element_name
                        in ["learning_objectives", "instructions", "safety"]
                        else "medium",
                        "suggestion": f"Missing {element_name.replace('_', ' ')}. This is important for effective {content_type} materials.",
                    }
                )

        # Check for pedagogical elements
        pedagogical_elements = {
            "engagement": [
                "activity",
                "interactive",
                "discuss",
                "participate",
                "engage",
            ],
            "assessment": ["check", "assess", "evaluate", "test your understanding"],
            "differentiation": [
                "extension",
                "challenge",
                "support",
                "scaffold",
                "adapt",
            ],
        }

        for element_name, keywords in pedagogical_elements.items():
            found = any(keyword in content_lower for keyword in keywords)
            if not found:
                gaps.append(
                    {
                        "element": element_name.title(),
                        "severity": "low",
                        "suggestion": f"Consider adding {element_name} elements to improve pedagogical effectiveness.",
                    }
                )

        return gaps

    async def batch_process(
        self, files: list[tuple[bytes, str, str | None]]
    ) -> list[dict[str, Any]]:
        """Process multiple files in batch"""
        results = []
        for file_content, filename, content_type in files:
            result = await self.process_file(file_content, filename, content_type)
            results.append(result)
        return results

    def categorize_content(self, content: str) -> dict[str, Any]:
        """Categorize content into appropriate course structure"""
        detected_type, confidence, type_scores = self._detect_content_type(content)
        sections = self._parse_sections(content)

        # Determine appropriate module/lesson placement
        return {
            "content_type": detected_type,
            "content_type_confidence": confidence,
            "alternative_types": [
                t
                for t, score in type_scores.items()
                if score > 0 and t != detected_type
            ][:3],
            "suggested_module": self._suggest_module(content, sections),
            "prerequisites": self._identify_prerequisites(content),
            "difficulty_level": self._assess_difficulty(content),
            "estimated_duration": self._estimate_duration(content, detected_type),
            "tags": self._extract_tags(content),
        }

    def _suggest_module(self, content: str, sections: list[dict]) -> str:
        """Suggest appropriate module based on content"""
        # Look for module/chapter indicators
        for section in sections:
            title = section.get("title", "").lower()
            if "module" in title or "chapter" in title or "unit" in title:
                return section["title"]

        # Default to content type
        return "General Module"

    def _identify_prerequisites(self, content: str) -> list[str]:
        """Identify prerequisite knowledge mentioned in content"""
        prerequisites = []
        patterns = [
            r"prerequisite[s]?:?\s*([^.]+)",
            r"prior knowledge:?\s*([^.]+)",
            r"you should know:?\s*([^.]+)",
            r"requires?:?\s*([^.]+)",
            r"assumes? familiarity with:?\s*([^.]+)",
        ]

        content_lower = content.lower()
        for pattern in patterns:
            matches = re.findall(pattern, content_lower)
            prerequisites.extend(matches)

        return list(set(prerequisites))[:5]  # Return up to 5 unique prerequisites

    def _assess_difficulty(self, content: str) -> str:
        """Assess difficulty level of content"""
        content_lower = content.lower()

        # Indicators of different difficulty levels
        beginner_indicators = [
            "basic",
            "introduction",
            "fundamental",
            "simple",
            "overview",
            "getting started",
        ]
        intermediate_indicators = [
            "intermediate",
            "application",
            "practice",
            "develop",
            "implement",
        ]
        advanced_indicators = [
            "advanced",
            "complex",
            "research",
            "theory",
            "optimization",
            "algorithm",
        ]

        beginner_score = sum(1 for ind in beginner_indicators if ind in content_lower)
        intermediate_score = sum(
            1 for ind in intermediate_indicators if ind in content_lower
        )
        advanced_score = sum(1 for ind in advanced_indicators if ind in content_lower)

        if advanced_score > max(beginner_score, intermediate_score):
            return "advanced"
        if intermediate_score > beginner_score:
            return "intermediate"
        return "beginner"

    def _estimate_duration(self, content: str, content_type: str) -> int:
        """Estimate duration in minutes based on content"""
        word_count = len(content.split())

        # Different reading/completion speeds for different content types
        speeds = {
            "lecture": 150,  # words per minute reading
            "quiz": 50,  # includes thinking time
            "worksheet": 75,  # includes problem solving
            "lab": 30,  # includes hands-on work
            "case_study": 100,  # includes analysis
            "general": 125,
        }

        speed = speeds.get(content_type, 125)
        duration = word_count // speed

        # Add minimum durations
        min_durations = {"quiz": 15, "worksheet": 20, "lab": 45, "case_study": 30}

        min_duration = min_durations.get(content_type, 10)
        return max(duration, min_duration)

    def _extract_tags(self, content: str) -> list[str]:
        """Extract relevant tags/keywords from content"""
        # Simple keyword extraction - in production, use NLP

        # Remove common words
        stop_words = {
            "the",
            "a",
            "an",
            "and",
            "or",
            "but",
            "in",
            "on",
            "at",
            "to",
            "for",
            "of",
            "with",
            "by",
            "from",
            "is",
            "are",
            "was",
            "were",
            "be",
            "been",
            "being",
            "have",
            "has",
            "had",
            "do",
            "does",
            "did",
            "will",
            "would",
            "could",
            "should",
            "may",
            "might",
            "must",
            "can",
            "this",
            "that",
            "these",
            "those",
            "i",
            "you",
            "he",
            "she",
            "it",
            "we",
            "they",
            "what",
            "which",
            "who",
            "when",
            "where",
            "why",
            "how",
            "all",
            "each",
            "every",
            "some",
            "any",
            "many",
            "much",
            "more",
            "most",
            "other",
            "another",
            "such",
            "no",
            "not",
            "only",
            "own",
            "same",
            "so",
            "than",
            "too",
            "very",
            "just",
        }

        # Extract words
        words = re.findall(r"\b[a-z]+\b", content.lower())

        # Filter and count
        filtered_words = [w for w in words if len(w) > 3 and w not in stop_words]
        word_counts = Counter(filtered_words)

        # Return top 10 most common words as tags
        return [word for word, _ in word_counts.most_common(10)]


# Create singleton instance
file_import_service = FileImportService()
